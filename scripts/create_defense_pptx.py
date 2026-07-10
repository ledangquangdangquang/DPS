from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "bao_ve_do_an_20_phut.pptx"
TOTAL_SLIDES = 18

NAVY = RGBColor(20, 46, 82)
BLUE = RGBColor(40, 95, 165)
LIGHT_BLUE = RGBColor(230, 238, 250)
GRAY = RGBColor(88, 96, 105)
LIGHT_GRAY = RGBColor(244, 246, 248)
BLACK = RGBColor(28, 31, 35)
WHITE = RGBColor(255, 255, 255)


def set_run(run, size=22, bold=False, color=BLACK):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_header(slide, title, number):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(33.867), Cm(1.18)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Cm(0.65), Cm(0.18), Cm(28.2), Cm(0.8))
    p = box.text_frame.paragraphs[0]
    p.margin_left = 0
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    set_run(run, 20, True, WHITE)

    num = slide.shapes.add_textbox(Cm(31.0), Cm(0.19), Cm(2.1), Cm(0.8))
    p = num.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{number:02d}/{TOTAL_SLIDES}"
    set_run(run, 13, False, WHITE)


def add_footer(slide):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(18.83), Cm(33.867), Cm(0.22)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE
    line.line.fill.background()


def add_title(slide, text, y=1.72, size=30):
    box = slide.shapes.add_textbox(Cm(0.85), Cm(y), Cm(32.0), Cm(1.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_run(run, size, True, NAVY)
    return box


def add_bullets(slide, items, x, y, w, h, size=21, color=BLACK):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(6)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_note(slide, text, x, y, w, h):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.color.rgb = BLUE
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Cm(0.35)
    tf.margin_right = Cm(0.35)
    tf.margin_top = Cm(0.25)
    tf.margin_bottom = Cm(0.2)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run(run, 18, True, NAVY)
    return shape


def add_agenda_card(slide, number, title, detail, x, y, w=14.8, h=2.05):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LIGHT_BLUE

    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Cm(x + 0.35), Cm(y + 0.38), Cm(1.25), Cm(1.25)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = BLUE
    badge.line.fill.background()
    tf = badge.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    set_run(run, 16, True, WHITE)

    box = slide.shapes.add_textbox(Cm(x + 1.95), Cm(y + 0.26), Cm(w - 2.35), Cm(h - 0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run(run, 17, True, NAVY)
    p = tf.add_paragraph()
    p.space_before = Pt(2)
    run = p.add_run()
    run.text = detail
    set_run(run, 12, False, GRAY)


def add_table(slide, data, x, y, w, h, font_size=16):
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, Cm(x), Cm(y), Cm(w), Cm(h))
    table = table_shape.table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = NAVY if r == 0 else (LIGHT_GRAY if r % 2 == 0 else WHITE)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    set_run(run, font_size, r == 0, WHITE if r == 0 else BLACK)
    return table_shape


def add_image(slide, rel_path, x, y, w, h):
    path = ROOT / rel_path
    if path.exists():
        slide.shapes.add_picture(str(path), Cm(x), Cm(y), width=Cm(w), height=Cm(h))
    else:
        add_note(slide, f"Thiếu hình: {rel_path}", x, y, w, h)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build():
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    # 1
    s = blank(prs)
    bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(33.867), Cm(19.05))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    add_title(
        s,
        "Mô phỏng kênh MIMO băng rộng biến thiên theo thời gian\n"
        "dựa trên mô hình hình học với độ phức tạp thấp bằng chuỗi DPS",
        y=3.0,
        size=31,
    )
    # recolor title text on cover
    for shape in s.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.font.color.rgb = WHITE
    add_bullets(
        s,
        [
            "Sinh viên thực hiện: ................................",
            "Giảng viên hướng dẫn: .............................",
            "Thời lượng trình bày: khoảng 20 phút",
        ],
        2.2,
        9.2,
        23.5,
        3.2,
        size=20,
        color=WHITE,
    )

    # 2
    s = blank(prs)
    add_header(s, "Mục lục trình bày", 2)
    add_title(s, "Lộ trình 20 phút", y=1.7, size=28)
    add_agenda_card(s, 1, "Đặt vấn đề", "Bối cảnh, SoCE và mục tiêu đồ án", 1.35, 3.35)
    add_agenda_card(s, 2, "Cơ sở phương pháp", "Giới hạn băng, DPS và mô hình MIMO 4D", 17.35, 3.35)
    add_agenda_card(s, 3, "Triển khai mô phỏng", "Các nhánh MATLAB, thuật toán DPS và cấu hình chính", 1.35, 6.25)
    add_agenda_card(s, 4, "Kết quả", "Đáp ứng kênh, NMSE, thời gian và benchmark", 17.35, 6.25)
    add_agenda_card(s, 5, "Kết luận", "Hạn chế, kết luận và hướng phát triển", 1.35, 9.15)
    add_note(s, "Mạch trình bày đi từ vấn đề độ phức tạp của SoCE đến cách dùng DPS và kiểm chứng bằng mô phỏng MATLAB.", 1.35, 13.0, 31.0, 2.05)
    add_footer(s)

    # 3
    s = blank(prs)
    add_header(s, "Bối cảnh nghiên cứu", 3)
    add_bullets(
        s,
        [
            "Kênh MIMO băng rộng thay đổi theo thời gian, tần số, anten phát và anten thu.",
            "Mô hình kênh dựa trên hình học (GCM) mô tả kênh bằng các thành phần đa đường (MPC).",
            "Mỗi MPC có trễ, Doppler, góc khởi hành, góc tới và trọng số phức.",
        ],
        1.2,
        2.3,
        16.1,
        7.0,
    )
    add_image(s, "results/figures/gcm_MIMO.png", 19.0, 2.2, 13.2, 8.0)
    add_note(s, "GCM có ý nghĩa vật lý rõ ràng, nhưng mô phỏng trực tiếp có thể tốn chi phí khi kích thước hệ thống tăng.", 1.2, 12.1, 31.0, 2.0)
    add_footer(s)

    # 4
    s = blank(prs)
    add_header(s, "Vấn đề của SoCE trực tiếp", 4)
    add_title(s, "Tổng các hàm mũ phức trên tensor kênh bốn chiều", y=1.65, size=24)
    add_note(s, "hₘ,q,s,r = Σ ηₚ exp{j2π(νₚm − θₚq + ζₚs − ξₚr)}", 1.2, 3.35, 31.1, 1.8)
    add_bullets(
        s,
        [
            "Mỗi MPC đóng góp một hàm mũ phức theo bốn chiều.",
            "SoCE phải cộng P MPC tại từng điểm mẫu của tensor kênh.",
            "Độ phức tạp: O(M Q NTx NRx P).",
            "Chi phí tăng nhanh khi tăng số mẫu, số anten hoặc số MPC.",
        ],
        1.45,
        6.25,
        29.7,
        6.5,
    )
    add_footer(s)

    # 5
    s = blank(prs)
    add_header(s, "Mục tiêu đồ án", 5)
    add_bullets(
        s,
        [
            "Hệ thống hóa mô hình GCM/SoCE cho kênh MIMO băng rộng biến thiên theo thời gian.",
            "Trình bày cơ sở DPS cho tín hiệu giới hạn băng trên khối mẫu hữu hạn.",
            "Triển khai MATLAB: SoCE, DPS chính xác, DPS xấp xỉ 4D và hybrid.",
            "Đánh giá sai số NMSE và thời gian tính toán của các nhánh mô phỏng.",
        ],
        1.25,
        2.3,
        30.4,
        8.4,
    )
    add_note(s, "Đồ án không tuyên bố đề xuất thuật toán DPS mới; trọng tâm là kiểm chứng và phân tích cách áp dụng DPS cho mô phỏng kênh.", 1.25, 12.3, 30.4, 2.25)
    add_footer(s)

    # 6
    s = blank(prs)
    add_header(s, "Ý tưởng chính", 6)
    add_bullets(
        s,
        [
            "Vận tốc cực đại giới hạn Doppler.",
            "Trễ cực đại giới hạn biến thiên theo tần số.",
            "Miền góc và khoảng cách anten giới hạn tần số không gian.",
            "Kênh có cấu trúc giới hạn băng theo bốn chiều.",
        ],
        1.25,
        2.4,
        15.5,
        7.2,
    )
    add_bullets(
        s,
        [
            "Kênh không phải tensor tùy ý.",
            "Có thể biểu diễn bằng không gian con có số chiều nhỏ hơn.",
            "DPS phù hợp với tín hiệu giới hạn băng trên đoạn chỉ số hữu hạn.",
        ],
        18.1,
        2.4,
        14.0,
        7.2,
    )
    add_note(s, "Thông điệp: dùng DPS để khai thác cấu trúc giới hạn băng, thay vì tính trực tiếp từng MPC tại từng mẫu.", 1.25, 12.1, 30.6, 2.1)
    add_footer(s)

    # 7
    s = blank(prs)
    add_header(s, "Cơ sở DPS", 7)
    add_bullets(
        s,
        [
            "DPS là các chuỗi rời rạc tập trung năng lượng tốt trong một miền tần số cho trước.",
            "Các vector DPS đầu tiên có trị riêng gần 1.",
            "Khi trị riêng giảm nhanh, số chiều thiết yếu nhỏ hơn nhiều so với số mẫu ban đầu.",
        ],
        1.2,
        2.2,
        16.0,
        7.0,
    )
    add_image(s, "results/figures/paper_figure4_dps_eigenvalues.png", 18.6, 2.1, 13.0, 8.0)
    add_note(s, "Biểu diễn không gian con: ĥᴰ = Vα", 1.2, 11.5, 14.4, 1.65)
    add_note(s, "DPS chuyển bài toán từ tính từng MPC tại từng mẫu sang biểu diễn trong không gian con.", 16.4, 11.5, 15.2, 1.65)
    add_footer(s)

    # 8
    s = blank(prs)
    add_header(s, "Mô hình MIMO băng rộng bốn chiều", 8)
    add_table(
        s,
        [
            ["Chiều", "Đại lượng vật lý", "Tần số chuẩn hóa"],
            ["Thời gian", "Doppler", "νₚ"],
            ["Tần số", "Trễ", "θₚ"],
            ["Anten phát", "AoD", "ζₚ"],
            ["Anten thu", "AoA", "ξₚ"],
        ],
        2.3,
        2.1,
        28.8,
        6.4,
        font_size=18,
    )
    add_note(s, "Mỗi chiều có một miền băng riêng; cơ sở DPS đa chiều được xây dựng từ các cơ sở một chiều.", 2.3, 10.5, 28.8, 2.0)
    add_footer(s)

    # 9
    s = blank(prs)
    add_header(s, "Các nhánh mô phỏng MATLAB", 9)
    add_table(
        s,
        [
            ["Nhánh", "Vai trò"],
            ["SoCE tham chiếu", "Tính trực tiếp từ MPC"],
            ["DPS chính xác", "Kiểm tra sai số do cắt không gian con"],
            ["DPS xấp xỉ 4D", "Xấp xỉ hệ số DPS trên cả bốn chiều"],
            ["Hybrid", "DPS theo thời gian--tần số, trực tiếp ở hai chiều anten"],
        ],
        1.2,
        2.0,
        31.2,
        6.5,
        font_size=15,
    )
    add_image(s, "results/figures/mimo_dps_kaltenberger_approx_flowchart.png", 5.8, 9.35, 22.2, 6.6)
    add_footer(s)

    # 10
    s = blank(prs)
    add_header(s, "Thuật toán DPS trong mô phỏng", 10)
    add_table(
        s,
        [
            ["Bước", "Nội dung", "Ý nghĩa"],
            ["1", "Tạo cơ sở DPS 1D cho thời gian, tần số, Tx, Rx", "Khai thác miền giới hạn băng từng chiều"],
            ["2", "Tính tensor hệ số α từ từng MPC", "Thay biểu diễn theo mẫu bằng biểu diễn theo hệ số"],
            ["3", "Tái tạo H bằng nhân tensor theo từng mode", "Sinh lại khối kênh MIMO bốn chiều"],
            ["4", "So sánh với SoCE bằng NMSE và thời gian", "Đánh giá sai số và chi phí tính toán"],
        ],
        1.1,
        2.0,
        31.5,
        6.3,
        font_size=13,
    )
    add_bullets(
        s,
        [
            "DPS chính xác: hệ số được tính bằng phép chiếu chính xác lên cơ sở đã cắt.",
            "DPS xấp xỉ 4D: hệ số được xấp xỉ từ tần số chuẩn hóa của MPC ở cả bốn chiều.",
            "Hybrid: chỉ xấp xỉ DPS ở thời gian--tần số, còn hai chiều anten dùng hàm mũ trực tiếp.",
        ],
        1.4,
        9.6,
        30.4,
        3.5,
        size=17,
    )
    add_note(s, "Điểm cần nói rõ: DPS không thay đổi tập MPC; phương pháp chỉ thay đổi cách biểu diễn và tái tạo cùng một kênh.", 1.4, 14.1, 30.4, 1.8)
    add_footer(s)

    # 11
    s = blank(prs)
    add_header(s, "Cấu hình mô phỏng chính", 11)
    add_table(
        s,
        [
            ["Tham số", "Giá trị", "Tham số", "Giá trị"],
            ["M", "256", "Q", "64"],
            ["NTx", "4", "NRx", "4"],
            ["P", "80", "Dt", "6"],
            ["Df", "9", "DTx = DRx", "4"],
            ["Dtotal", "864", "Số mẫu kênh", "262144"],
        ],
        2.0,
        2.25,
        29.8,
        7.2,
        font_size=17,
    )
    add_note(s, "Số hệ số DPS nhỏ hơn đáng kể số phần tử kênh, nhưng thời gian chạy không giảm đúng theo tỷ lệ này.", 2.0, 11.5, 29.8, 2.0)
    add_footer(s)

    # 12
    s = blank(prs)
    add_header(s, "Kết quả minh họa đáp ứng kênh", 12)
    add_image(s, "results/figures/mimo_dps_kaltenberger_approx_tx1_rx1.png", 1.0, 1.55, 31.8, 12.6)
    add_note(s, "Hình chỉ minh họa một cặp anten; kết luận định lượng dựa trên NMSE tính trên toàn bộ tensor.", 2.3, 15.0, 29.0, 1.55)
    add_footer(s)

    # 13
    s = blank(prs)
    add_header(s, "Chỉ tiêu sai số NMSE", 13)
    add_note(s, "NMSE = E{|Href − Htest|²} / E{|Href|²}", 3.0, 2.35, 27.9, 1.8)
    add_bullets(
        s,
        [
            "Href: kênh tham chiếu SoCE.",
            "Htest: kênh tái tạo từ nhánh DPS.",
            "NMSE càng nhỏ thì sai số tương đối càng thấp.",
            "Phép trung bình được tính trên toàn bộ tensor kênh.",
        ],
        2.0,
        5.35,
        29.6,
        6.0,
    )
    add_note(s, "NMSE giúp so sánh chất lượng tái tạo giữa các nhánh trên cùng cấu hình mô phỏng.", 2.0, 12.8, 29.6, 1.9)
    add_footer(s)

    # 14
    s = blank(prs)
    add_header(s, "Kết quả NMSE", 14)
    add_image(s, "results/figures/mimo_dps_kaltenberger_approx_nmse_bar.png", 1.0, 1.7, 14.8, 10.8)
    add_table(
        s,
        [
            ["Nhánh", "NMSE"],
            ["DPS chính xác", "8,53×10⁻⁹"],
            ["DPS xấp xỉ 4D", "1,99×10⁻⁴"],
            ["Hybrid", "4,29×10⁻⁷"],
            ["Hybrid so với DPS chính xác", "4,20×10⁻⁷"],
        ],
        16.6,
        2.25,
        15.6,
        6.1,
        font_size=14,
    )
    add_bullets(
        s,
        [
            "DPS chính xác: sai số cắt rất nhỏ trong cấu hình hiện tại.",
            "DPS xấp xỉ 4D: sai số lớn hơn do xấp xỉ hệ số ở cả bốn chiều.",
            "Hybrid: giảm sai số vì tránh xấp xỉ ở hai chiều anten.",
        ],
        16.6,
        9.35,
        15.6,
        4.8,
        size=17,
    )
    add_footer(s)

    # 15
    s = blank(prs)
    add_header(s, "Kết quả thời gian tính toán", 15)
    add_image(s, "results/figures/mimo_dps_kaltenberger_approx_runtime_bar.png", 1.0, 1.7, 14.8, 10.8)
    add_table(
        s,
        [
            ["Nhánh", "Thời gian"],
            ["SoCE", "0,222372 s"],
            ["DPS chính xác", "0,045555 s"],
            ["DPS xấp xỉ 4D", "0,080983 s"],
            ["Hybrid", "0,025322 s"],
        ],
        16.7,
        2.35,
        15.2,
        5.8,
        font_size=15,
    )
    add_note(s, "Trong cấu hình hiện tại, hybrid có thời gian xử lý nhỏ nhất. Phép đo chưa bao gồm đầy đủ chi phí tạo cơ sở DPS.", 16.7, 9.4, 15.2, 2.6)
    add_footer(s)

    # 16
    s = blank(prs)
    add_header(s, "Kết quả bổ sung", 16)
    add_image(s, "results/figures/mimo_dps_kaltenberger_p_sweep_runtime.png", 1.0, 1.7, 14.8, 9.8)
    add_image(s, "results/figures/mimo_four_method_benchmark_nmse.png", 17.3, 1.7, 14.8, 9.8)
    add_bullets(
        s,
        [
            "Sweep theo số MPC: P = 10, 20, 40, 80, 160, 320.",
            "Benchmark: SoCE trực tiếp, SoCE đệ quy, CE-BEM và DPS chính xác.",
            "Không kết luận DPS luôn tốt hơn; kết quả phụ thuộc cấu hình, miền băng và cách triển khai.",
        ],
        1.2,
        12.35,
        31.0,
        3.4,
        size=17,
    )
    add_footer(s)

    # 17
    s = blank(prs)
    add_header(s, "Hạn chế của đồ án", 17)
    add_bullets(
        s,
        [
            "Chưa tái tạo đầy đủ một hình hoặc bảng cụ thể của bài báo gốc với toàn bộ tham số khớp nhau.",
            "Cấu hình mô phỏng chính còn hẹp.",
            "Số chiều DPS và hệ số phân giải chưa được tối ưu tự động trong nhánh chính.",
            "Thời gian chạy phụ thuộc môi trường MATLAB, cách vector hóa và phần cứng.",
            "Chi phí tạo cơ sở DPS chưa được tính đầy đủ trong phép đo thời gian chính.",
        ],
        1.4,
        2.1,
        30.5,
        8.5,
    )
    add_note(s, "Kết quả cần được hiểu trong phạm vi cấu hình mô phỏng hiện tại.", 1.4, 12.35, 30.5, 2.0)
    add_footer(s)

    # 18
    s = blank(prs)
    add_header(s, "Kết luận và hướng phát triển", 18)
    add_bullets(
        s,
        [
            "Đã trình bày mô hình GCM/SoCE cho kênh MIMO băng rộng biến thiên theo thời gian.",
            "Đã triển khai và so sánh SoCE, DPS chính xác, DPS xấp xỉ 4D và hybrid.",
            "Trong cấu hình hiện tại, DPS chính xác có sai số cắt nhỏ.",
            "Hybrid đạt cân bằng tốt hơn giữa sai số và thời gian so với xấp xỉ 4D.",
        ],
        1.3,
        2.1,
        30.6,
        5.7,
        size=19,
    )
    add_bullets(
        s,
        [
            "Hướng phát triển: quét thêm số anten, số MPC, miền băng và số chiều DPS.",
            "Chọn số chiều DPS và hệ số phân giải theo ngưỡng sai số.",
            "Tái tạo đầy đủ một kết quả định lượng cụ thể của bài báo Kaltenberger.",
            "Mở rộng sang mô hình kênh chuẩn thực tế hơn.",
        ],
        1.3,
        8.55,
        30.6,
        5.5,
        size=19,
    )
    add_note(s, "Em xin kết thúc phần trình bày và kính mong nhận được góp ý của quý thầy cô.", 3.4, 15.25, 27.0, 1.7)
    add_footer(s)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
